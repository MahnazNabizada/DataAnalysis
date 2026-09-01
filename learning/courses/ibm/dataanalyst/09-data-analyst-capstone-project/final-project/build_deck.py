"""Build "Data Analyst Capstone Project Report.pptx" from the course template.

The template is opened in place (as a copy) so every slide keeps the IBM Skills
Network master, layouts, theme fonts and decorative artwork. Instruction text is
replaced with the real analysis, the "<chart goes here>" boxes are swapped for
the PNGs from build_charts.py, and extra result/appendix slides are appended and
then re-ordered into position.
"""

import copy
import json
import os
import shutil

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.shapes.autoshape import CT_Shape
from pptx.shapes.autoshape import Shape
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(HERE, "Data Analyst Capstone Template 2026.pptx")
ASSETS = os.path.join(HERE, "assets")
OUT_PPTX = os.path.join(HERE, "Data Analyst Capstone Project Report.pptx")

AUTHOR = "Pablo Fiterman"
DATE = "September 1, 2026"
WATERMARK = f"{AUTHOR} · IBM Data Analyst Capstone · {DATE}"

INK = RGBColor(0x26, 0x26, 0x26)
INK2 = RGBColor(0x52, 0x52, 0x52)

facts = json.load(open(os.path.join(ASSETS, "facts.json"), encoding="utf-8"))

# ============================================================ xml/text helpers


def drop(shape):
    shape._element.getparent().remove(shape._element)


def _clone_pPr(src_p, level):
    """Copy a source paragraph's properties so bullets/indent survive a rewrite."""
    pPr = src_p._p.find(qn("a:pPr"))
    new = copy.deepcopy(pPr) if pPr is not None else None
    if new is not None:
        if level:
            new.set("lvl", str(level))
        else:
            new.attrib.pop("lvl", None)
    return new


def _clone_rPr(src_p):
    for r in src_p.runs:
        rPr = r._r.find(qn("a:rPr"))
        if rPr is not None:
            return copy.deepcopy(rPr)
    return None


def write(tf, items, base=0, size=None, tight=False):
    """Rewrite a text frame from `items` = [(text, level), ...] or [text, ...].

    Formatting (bullet glyph, indent, font, colour) is inherited from the
    template paragraph at index `base`, so rewritten slides still look native.
    Pass an int for `size` to force a point size, or a dict {level: pt}.
    """
    src = tf.paragraphs[min(base, len(tf.paragraphs) - 1)]
    pPr_tpl = {}
    rPr_tpl = {}
    for p in tf.paragraphs:
        lvl = p.level
        if lvl not in pPr_tpl:
            pPr_tpl[lvl] = _clone_pPr(p, lvl)
            rPr_tpl[lvl] = _clone_rPr(p)
    fallback_pPr, fallback_rPr = _clone_pPr(src, 0), _clone_rPr(src)

    body = tf._txBody
    for p in list(body.findall(qn("a:p"))):
        body.remove(p)

    norm = [(i, 0) if isinstance(i, str) else i for i in items]
    for text, level in norm:
        p = body.makeelement(qn("a:p"), {})
        body.append(p)
        pPr = pPr_tpl.get(level)
        if pPr is None:
            pPr = fallback_pPr
        if pPr is not None:
            new_pPr = copy.deepcopy(pPr)
            if level:
                new_pPr.set("lvl", str(level))
            else:
                new_pPr.attrib.pop("lvl", None)
            p.append(new_pPr)
        r = body.makeelement(qn("a:r"), {})
        p.append(r)
        rPr = rPr_tpl.get(level)
        if rPr is None:
            rPr = fallback_rPr
        if rPr is not None:
            r.append(copy.deepcopy(rPr))
        t = body.makeelement(qn("a:t"), {})
        t.text = text
        r.append(t)

    # Point sizes: dict keyed by level, or one size for everything.
    for p in tf.paragraphs:
        want = size.get(p.level) if isinstance(size, dict) else size
        if want:
            for r in p.runs:
                r.font.size = Pt(want)
    if tight:
        tf.word_wrap = True
    return tf


def fit(tf, shrink=True):
    """Ask PowerPoint to shrink text on overflow (normAutofit)."""
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        return
    for tag in ("a:normAutofit", "a:spAutoFit", "a:noAutofit"):
        el = bodyPr.find(qn(tag))
        if el is not None:
            bodyPr.remove(el)
    if shrink:
        bodyPr.append(bodyPr.makeelement(qn("a:normAutofit"), {}))


def picture(slide, name, left, top, max_w, max_h):
    """Place a PNG centred inside the (left, top, max_w, max_h) box, aspect kept."""
    path = os.path.join(ASSETS, name)
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    w, h = iw * scale, ih * scale
    return slide.shapes.add_picture(
        path,
        Inches(left + (max_w - w) / 2),
        Inches(top + (max_h - h) / 2),
        width=Inches(w),
        height=Inches(h),
    )


def note(slide, text, left, top, width, size=11, color=INK2, align=PP_ALIGN.LEFT,
         bold=False, italic=True, height=0.42):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.italic = italic
    r.font.bold = bold
    return box


def find(slide, predicate):
    return [sh for sh in slide.shapes if predicate(sh)]


def by_text(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle.lower() in sh.text_frame.text.lower():
            return sh
    return None


# ============================================================ deck construction
shutil.copyfile(TEMPLATE, OUT_PPTX)
prs = Presentation(OUT_PPTX)
S = prs.slides
BLANK = prs.slide_layouts[5]      # "2_Blank" - no placeholders, keeps the master art

# ---------------------------------------------------------- watermark on master
# The brief allows a master-level text box used as a watermark on every slide.
# MasterShapes has no add_textbox(), so the sp element is built directly and
# appended to the master's shape tree — and to each layout as well, since a
# layout that opts out of inheriting master shapes would otherwise drop it.
def add_watermark(shape_collection, text):
    spTree = shape_collection._spTree
    used = {int(el.get("id")) for el in spTree.iter() if el.tag.endswith("}cNvPr")
            and el.get("id") and el.get("id").isdigit()}
    sp = CT_Shape.new_textbox_sp(
        max(used, default=1) + 1, "Watermark",
        Inches(0.28), Inches(7.05), Inches(12.8), Inches(0.32))
    spTree.append(sp)
    box = Shape(sp, None)
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(9)
    r.font.color.rgb = RGBColor(0xA8, 0xA8, 0xA8)
    r.font.italic = True
    return box


add_watermark(prs.slide_master.shapes, WATERMARK)
for layout in prs.slide_layouts:
    add_watermark(layout.shapes, WATERMARK)

# ------------------------------------------------------------------ 1. TITLE
s = S[0]
write(s.shapes[0].text_frame, ["DATA ANALYST CAPSTONE PROJECT"])
write(s.shapes[1].text_frame, [
    "Emerging Technology Trends in the Developer Job Market",
    "",
    AUTHOR,
    DATE,
    "",
    "IBM Data Analyst Professional Certificate",
    "Stack Overflow Developer Survey · Jobs API · Web-scraped salary data",
], size={0: 18})
for i, para in enumerate(s.shapes[1].text_frame.paragraphs):
    for r in para.runs:
        r.font.size = Pt(20 if i == 0 else 15)
        r.font.bold = i == 0
# Template subtitle box runs under the decorative image on the right; trim it.
s.shapes[1].width = Inches(5.15)
fit(s.shapes[1].text_frame)

# ------------------------------------------------------------------ 2. OUTLINE
s = S[1]
write(by_text(s, "Executive Summary").text_frame, [
    "Executive Summary",
    "Introduction",
    "Methodology",
    "Results — job market, languages, databases",
    "Dashboard — current usage, future trends, demographics",
    "Discussion",
    "Findings & Implications",
    "Conclusion",
    "Appendix",
], size=17)

# --------------------------------------------------------- 3. EXECUTIVE SUMMARY
s = S[2]
write(by_text(s, "Summarize key findings").text_frame, [
    ("Three data sources were combined to read the developer job market: "
     f"{facts['rows']:,} Stack Overflow survey responses, "
     f"{facts['job_total']:,} job postings from a Jobs API and scraped salary data "
     "for 10 languages.", 0),
    ("Demand is concentrated. Three metro areas hold "
     f"{facts['job_top3_share']}% of all postings, led by "
     f"{facts['job_top'][0]} with {facts['job_top'][1]:,}.", 0),
    ("The current stack is stable: JavaScript, SQL and HTML/CSS lead usage, "
     "PostgreSQL leads databases, AWS leads platforms.", 0),
    ("Intent diverges from usage. Rust (+3.3 pp) and Go (+2.5 pp) gain the most "
     "mind-share for next year; JavaScript, PHP and Java lose the most, while "
     "still leading in absolute terms.", 0),
    ("PostgreSQL is the only database in the current top five that grows in "
     f"absolute terms ({facts['Database_have_top5'][0][1]:,} to "
     f"{facts['Database_want_top5'][0][1]:,}); Redis gains from sixth place.", 0),
    (f"Pay follows scarcity, not popularity: {facts['sal_top'][0]} tops the salary "
     f"table at ${facts['sal_top'][1]:,} while the most-used languages sit mid-table.", 0),
    (f"Compensation correlates only weakly with experience (r = "
     f"{facts['corr_pairs']['comp_workexp']}) and barely at all with satisfaction "
     f"(r = {facts['corr_pairs']['jobsat_comp']}); country explains far more.", 0),
], size=13)
fit(by_text(s, "Three data sources").text_frame)

# --------------------------------------------------------------- 4. INTRODUCTION
s = S[3]
write(by_text(s, "Purpose of the report").text_frame, [
    ("Purpose, audience and value", 0),
    ("Purpose — establish which technologies developers use today, which they "
     "intend to adopt next year, and where the job market pays for them.", 1),
    ("Audience — hiring managers, L&D leads and technology strategists planning "
     "headcount, training budgets and platform bets for the next cycle.", 1),
    ("Business questions", 0),
    ("Where are the job openings, and for which skills?", 1),
    ("Which languages, databases, platforms and frameworks dominate today?", 1),
    ("Which are gaining or losing developer mind-share for next year?", 1),
    ("Who are the developers behind these answers, and what do they earn?", 1),
    ("Value — the current-versus-desired gap is a leading indicator: it flags "
     "skills to hire for before the market reprices them.", 0),
], size={0: 15, 1: 13})
fit(by_text(s, "Purpose, audience").text_frame)

# ---------------------------------------------------------------- 5. METHODOLOGY
s = S[4]
write(by_text(s, "Data sources").text_frame, [
    ("Data sources", 0),
    (f"Jobs API (REST/JSON) — {facts['job_total']:,} postings across "
     f"{facts['job_n_loc']} US metro areas, plus counts per technology.", 1),
    ("Web scraping (requests + BeautifulSoup) — annual average salary for 10 "
     "popular programming languages.", 1),
    (f"Stack Overflow Developer Survey — {facts['rows']:,} responses, "
     f"{facts['cols']} columns, {facts['countries']} countries.", 1),
    ("Data wrangling", 0),
    ("Duplicates checked and removed; column names normalised.", 1),
    ("Missing values treated per column: median for skewed numerics, mode for "
     "categoricals, and left blank where over half the column was absent.", 1),
    ("Min-max scaling and a log1p transform applied to compensation.", 1),
    ("Multi-select answers split on ';' and exploded to one row per selection.", 1),
    ("Analysis & tools", 0),
    ("Python (pandas, NumPy), SQLite/SQL, Matplotlib & Seaborn, Google Looker "
     "Studio for the dashboard.", 1),
], size={0: 14, 1: 12})
fit(by_text(s, "Data sources").text_frame)

# =============================== new slides appended, re-ordered at the end ====
new_positions = []  # (slide, target_index_after_reorder)


def add_blank(title=None, title_size=26):
    slide = prs.slides.add_slide(BLANK)
    if title:
        box = slide.shapes.add_textbox(Inches(0.62), Inches(0.30), Inches(12.1), Inches(0.72))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(title_size)
        r.font.bold = True
        r.font.color.rgb = INK
    return slide


def bullets(slide, items, left, top, width, height, size):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for text, level in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        r = p.add_run()
        r.text = ("" if level == 0 else "•  ") + text
        r.font.size = Pt(size if level == 0 else size - 1)
        r.font.bold = level == 0 and text.endswith(":")
        r.font.color.rgb = INK if level == 0 else INK2
    fit(tf)
    return box


# ---- 6. METHODOLOGY — DATA QUALITY (new, sits right after Methodology)
s6 = add_blank("METHODOLOGY — DATA QUALITY & TREATMENT")
picture(s6, "a10_missing_values.png", 0.55, 1.15, 6.6, 5.4)
bullets(s6, [
    ("What the profiling found:", 0),
    ("No duplicate rows and no duplicate ResponseIds remain after the "
     "de-duplication step, so the row count is trustworthy.", 1),
    (f"ConvertedCompYearly is missing for {facts['comp_missing_pct']}% of rows — "
     "the single biggest constraint on any pay analysis.", 1),
    (f"JobSat ({facts['jobsat_missing_pct']}%) and WorkExp "
     f"({facts['miss_top'][2][1]}%) are also heavily incomplete.", 1),
    ("How it was handled:", 0),
    ("Compensation was left unimputed for distribution and correlation work and "
     "reported on its valid subset "
     f"({facts['comp_valid']:,} responses) — imputing half a column with its own "
     "median would manufacture a false peak.", 1),
    ("Outliers were bounded with the 1.5 × IQR rule; "
     f"{facts['comp_outliers']:,} values ({facts['comp_outlier_pct']}%) fall above "
     f"${facts['comp_iqr_upper']:,} and are excluded from pay comparisons, not deleted "
     "from the dataset.", 1),
    ("Every count on the following slides therefore states its own denominator.", 0),
], 7.35, 1.15, 5.4, 5.4, 13)
new_positions.append((s6, 5))

# --------------------------------------------------- 7. RESULTS section divider
s = S[5]
write(s.placeholders[0].text_frame, ["RESULTS"])  # template ships it lower-case
tb = by_text(s, "Please present your results")
write(tb.text_frame, [
    "Job market · Programming languages · Databases",
], size=20)
tb.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
# The empty content box on this slide gets the section contents list.
empty = [sh for sh in s.shapes
         if sh.has_text_frame and not sh.text_frame.text.strip() and sh != tb]
if empty:
    write(empty[0].text_frame, [
        ("Job postings by location", 0),
        ("Popular languages by salary", 0),
        ("Programming language trends", 0),
        ("Database trends", 0),
    ], size=17)
    fit(empty[0].text_frame)

# -------------------------------------------------------------- 8. JOB POSTINGS
s = S[6]
write(s.placeholders[0].text_frame, ["JOB POSTINGS"])  # template ships a leading space
tb = by_text(s, "In Module 1 you have collected the job posting data")
tb.left, tb.top, tb.width, tb.height = (Inches(7.25), Inches(1.55),
                                        Inches(5.45), Inches(4.9))
write(tb.text_frame, [
    (f"{facts['job_total']:,} Python postings across "
     f"{facts['job_n_loc']} metro areas", 0),
    (f"{facts['job_top'][0]} leads with {facts['job_top'][1]:,} postings — "
     f"{round(100 * facts['job_top'][1] / facts['job_total'], 1)}% of the total.", 0),
    (f"The top three metros hold {facts['job_top3_share']}% of all postings; "
     "demand is geographically concentrated rather than evenly spread.", 0),
    (f"The tail is very thin: {facts['job_bottom'][0]} returns only "
     f"{facts['job_bottom'][1]} postings, a "
     f"{round(facts['job_top'][1] / facts['job_bottom'][1]):,}× gap to the leader.", 0),
    ("Read as relative demand, not absolute market size — the API returns "
     "postings indexed at collection time, not a census of open roles.", 0),
], size=13)
fit(tb.text_frame)
picture(s, "r1_job_postings.png", 0.55, 1.45, 6.5, 5.1)

# ---------------------------------------------------------- 9. POPULAR LANGUAGES
s = S[7]
tb = by_text(s, "In Module 1 you have collected the job postings data using web scraping")
tb.left, tb.top, tb.width, tb.height = (Inches(7.25), Inches(1.55),
                                        Inches(5.45), Inches(4.9))
write(tb.text_frame, [
    ("Salary rewards scarcity, not ubiquity", 0),
    (f"{facts['sal_top'][0]} pays the most at ${facts['sal_top'][1]:,}, ahead of "
     f"Python (${facts['sal_python']:,}) and C++.", 0),
    (f"{facts['sal_bottom'][0]} sits last at ${facts['sal_bottom'][1]:,} — a "
     f"${facts['sal_spread']:,} spread across the ten languages "
     f"(mean ${facts['sal_mean']:,}).", 0),
    ("The most widely used languages are not the best paid: JavaScript and SQL "
     "dominate survey usage yet sit in the lower half of this table.", 0),
    ("Implication — specialising in a narrower, in-demand language is better "
     "compensated than adding another mainstream one.", 0),
], size=13)
fit(tb.text_frame)
picture(s, "r2_language_salary.png", 0.55, 1.45, 6.5, 5.1)

# ------------------------------------------------- 10. PROGRAMMING LANG TRENDS
s = S[8]
for sh in find(s, lambda x: x.has_text_frame and "Bar chart of top 10 programming" in x.text_frame.text):
    drop(sh)
tb = by_text(s, "Summarize key trends shown in the charts")
tb.left, tb.top, tb.width, tb.height = (Inches(0.72), Inches(1.32),
                                        Inches(11.9), Inches(0.85))
write(tb.text_frame, [
    "The top five are unchanged in composition, but four of them lose share while "
    "TypeScript gains and climbs past HTML/CSS; Go and Rust break into the top ten.",
], size=13)
fit(tb.text_frame)
picture(s, "r3_lang_now.png", 0.60, 2.55, 5.85, 4.1)
picture(s, "r3_lang_next.png", 6.75, 2.55, 5.85, 4.1)

# ------------------------ 11. PROGRAMMING LANGUAGE TRENDS — FINDINGS/IMPLICATIONS
s = S[9]
ph = {p.placeholder_format.idx: p for p in s.placeholders}
write(ph[1].text_frame, [
    ("Findings", 0),
    ("", 0),
    (f"JavaScript leads both years but loses the most share of any language "
     f"({facts['Language_losers'][0][1]} pp).", 0),
    (f"Rust is the biggest gainer at +{facts['Language_gainers'][0][1]} pp, "
     f"followed by Go at +{facts['Language_gainers'][1][1]} pp.", 0),
    ("TypeScript is the only current top-four language that gains share "
     f"(+{[v for k, v in facts['Language_gainers'] if k == 'TypeScript'][0]} pp), "
     "overtaking HTML/CSS for next year.", 0),
    (f"PHP ({facts['Language_losers'][1][1]} pp) and Java "
     f"({facts['Language_losers'][2][1]} pp) fall furthest after JavaScript.", 0),
    ("Respondents list fewer languages they want "
     f"({facts['Language_per_resp'][1]} on average) than they already use "
     f"({facts['Language_per_resp'][0]}) — intent is narrower than practice.", 0),
], size=12)
write(ph[2].text_frame, [
    ("Implications", 0),
    ("", 0),
    ("Keep hiring for the incumbents — JavaScript, SQL and TypeScript remain the "
     "baseline, and losing share is not the same as losing volume.", 0),
    ("Treat Rust and Go as a build-ahead investment: demand for them is being "
     "created by developer preference before the job market reprices it.", 0),
    ("Fund TypeScript training explicitly; it is the only mainstream language "
     "gaining share and it consolidates the JavaScript ecosystem.", 0),
    ("Plan PHP and Java estates as maintenance, not growth — expect a shrinking "
     "and more expensive hiring pool over time.", 0),
    ("Because intent lists are shorter than usage lists, read these as changes in "
     "priority, not as forecasts of abandonment.", 0),
], size=12)
for t in (ph[1].text_frame, ph[2].text_frame):
    fit(t)

# ------------------------------------------------------------ 12. DATABASE TRENDS
s = S[10]
for sh in find(s, lambda x: x.has_text_frame and "Bar chart of top 10 databases" in x.text_frame.text):
    drop(sh)
tb = by_text(s, "Summarize key trends shown in the charts")
tb.left, tb.top, tb.width, tb.height = (Inches(0.72), Inches(1.32),
                                        Inches(11.9), Inches(0.85))
write(tb.text_frame, [
    "PostgreSQL extends its lead and is the only database in the current top five "
    "to grow; Redis jumps from sixth to second while MySQL drops to fourth.",
], size=13)
fit(tb.text_frame)
picture(s, "r3_db_now.png", 0.60, 2.55, 5.85, 4.1)
picture(s, "r3_db_next.png", 6.75, 2.55, 5.85, 4.1)

# --------------------------- 13. DATABASE TRENDS — FINDINGS & IMPLICATIONS
s = S[11]
ph = {p.placeholder_format.idx: p for p in s.placeholders}
write(ph[1].text_frame, [
    ("Findings", 0),
    ("", 0),
    (f"PostgreSQL leads both years and grows in absolute terms — "
     f"{facts['Database_have_top5'][0][1]:,} current users versus "
     f"{facts['Database_want_top5'][0][1]:,} intending to use it next year "
     f"(+{facts['Database_gainers'][0][1]} pp share).", 0),
    (f"MySQL is the largest loser at {facts['Database_losers'][0][1]} pp, falling "
     "from second place to fourth.", 0),
    (f"Redis gains +{facts['Database_gainers'][1][1]} pp and moves from sixth to "
     "second — caching and in-memory workloads are becoming default.", 0),
    (f"Supabase enters the top ten (+{facts['Database_gainers'][2][1]} pp) with no "
     "comparable current-usage footprint.", 0),
    ("Commercial engines decline: Microsoft SQL Server "
     f"{facts['Database_losers'][1][1]} pp and Oracle "
     f"{facts['Database_losers'][3][1]} pp.", 0),
], size=12)
write(ph[2].text_frame, [
    ("Implications", 0),
    ("", 0),
    ("Standardise new services on PostgreSQL — it is the only clear consolidation "
     "point, which lowers both hiring and operating risk.", 0),
    ("Budget for Redis as core infrastructure rather than an optimisation, and "
     "hire for cache-aware design.", 0),
    ("Plan MySQL, SQL Server and Oracle estates for migration or long-term "
     "maintenance; the skills pool is drifting away from them.", 0),
    ("Watch Supabase and BigQuery as managed-service entrants: developer interest "
     "arrives before enterprise adoption.", 0),
    ("Interest is not deployment — validate these signals against your own "
     "workload, licensing and data-residency constraints before committing.", 0),
], size=12)
for t in (ph[1].text_frame, ph[2].text_frame):
    fit(t)

# ------------------------------------------------ 14. DASHBOARD section divider
s = S[12]
tb = by_text(s, "Please present your dashboard")
tb.left, tb.top, tb.width, tb.height = (Inches(5.0), Inches(2.0),
                                        Inches(7.5), Inches(4.4))
write(tb.text_frame, [
    ("Three tabs built on the survey extract", 0),
    ("Tab 1 — Current technology usage: top 10 languages, databases, platforms "
     "and web frameworks.", 0),
    ("Tab 2 — Future technology trends: the same four dimensions for next year.", 0),
    ("Tab 3 — Demographics: respondents by country, age group and education "
     "level.", 0),
    (f"Source: one extract of {facts['rows']:,} survey responses, with each "
     "multi-select field normalised to one row per selection.", 0),
], size=14)
fit(tb.text_frame)

# ------------------------------------------------------------- 15-17. DASHBOARD
dash = [
    (13, "d1_current_usage.png", [
        ("Dashboard insight — the working stack is web-first and cloud-centred", 0),
        (f"JavaScript ({facts['Language_have_top5'][0][1]:,}), SQL "
         f"({facts['Language_have_top5'][1][1]:,}) and HTML/CSS "
         f"({facts['Language_have_top5'][2][1]:,}) are used by the largest share of "
         "respondents; TypeScript is already fourth.", 1),
        (f"PostgreSQL ({facts['Database_have_top5'][0][1]:,}) is the most-used "
         f"database, ahead of MySQL ({facts['Database_have_top5'][1][1]:,}); SQLite "
         "third place reflects embedded and local development.", 1),
        (f"AWS ({facts['Platform_have_top5'][0][1]:,}) is used by more respondents "
         f"than Azure ({facts['Platform_have_top5'][1][1]:,}) and Google Cloud "
         f"({facts['Platform_have_top5'][2][1]:,}) individually, and nearly as many as "
         "the two combined.", 1),
        (f"Node.js ({facts['Webframe_have_top5'][0][1]:,}) and React "
         f"({facts['Webframe_have_top5'][1][1]:,}) dominate frameworks, but jQuery "
         f"is still third ({facts['Webframe_have_top5'][2][1]:,}) — a large legacy "
         "estate is still in production.", 1),
    ]),
    (14, "d2_future_trends.png", [
        ("Dashboard insight — intent shifts toward typed, systems and managed tools", 0),
        ("The language leaders hold their ranks, but Go (5,661) and Rust (5,597) "
         "enter the top ten at sixth and seventh, and TypeScript overtakes "
         "HTML/CSS for third.", 1),
        (f"PostgreSQL grows to {facts['Database_want_top5'][0][1]:,} and Redis rises "
         f"to second ({facts['Database_want_top5'][1][1]:,}); MySQL falls to fourth.", 1),
        ("The cloud hierarchy is stable — AWS, Azure, Google Cloud — but Cloudflare "
         f"gains ({facts['Platform_have_top5'][3][1]:,} → "
         f"{facts['Platform_want_top5'][3][1]:,}) and Heroku drops out of the top ten.", 1),
        ("React overtakes Node.js in frameworks and Next.js moves to third, while "
         "jQuery — currently third by usage — drops out of the desired "
         "top ten.", 1),
    ]),
    (15, "d3_demographics.png", [
        ("Dashboard insight — a young, formally educated, Western-weighted sample", 0),
        (f"The top five countries supply {facts['country_top5_share']}% of all "
         f"responses, led by the United States ({facts['country_top5'][0][1]:,}); "
         f"{facts['countries']} countries appear in total.", 1),
        (f"The 25–34 bracket alone is {facts['age_2534_share']}% of respondents and "
         f"{facts['age_under35_share']}% are under 35 — the sample skews early-career.", 1),
        (f"{facts['ed_degree_share']}% hold a tertiary degree, with "
         f"{facts['ed_top'][0]} the single largest group "
         f"({facts['ed_top'][1]:,} respondents).", 1),
        ("Caveat — this is a self-selected sample of Stack Overflow users, so it "
         "over-represents English-speaking, web-focused developers. Every trend above "
         "should be read as this population's direction of travel, not the whole "
         "industry's.", 1),
    ]),
]
for idx, img, insight in dash:
    s = S[idx]
    ph = {p.placeholder_format.idx: p for p in s.placeholders}
    drop(ph[1])
    # The three template title boxes are 1.08-1.45in tall at differing tops;
    # normalise them so the wide dashboard image clears the title on all three.
    title = ph[0]
    title.left, title.top = Inches(0.66), Inches(0.30)
    title.width, title.height = Inches(12.1), Inches(0.86)
    fit(title.text_frame)
    picture(s, img, 0.45, 1.24, 12.45, 4.18)
    bullets(s, insight, 0.62, 5.48, 12.1, 1.50, 12)

# ---------------------------------------------------------------- 18. DISCUSSION
s = S[16]
ph = {p.placeholder_format.idx: p for p in s.placeholders}
drop(ph[1])
picture(s, "a11_lang_have_vs_want.png", 0.45, 1.55, 6.1, 5.0)
tb = ph[2]
tb.left, tb.top, tb.width, tb.height = (Inches(6.85), Inches(1.40),
                                        Inches(5.85), Inches(5.35))
write(tb.text_frame, [
    ("What the three tabs say together", 0),
    ("Continuity, not disruption. The same technologies lead current usage and "
     "next-year intent; what changes is the ordering and the size of the gaps.", 0),
    ("Almost every incumbent loses share, partly because respondents name fewer "
     "technologies they want than they already use. Comparing shares rather than "
     "raw counts is what separates a real shift from that arithmetic.", 0),
    ("Consolidation is the dominant pattern: PostgreSQL in databases, "
     "TypeScript in languages, React/Next.js in frameworks, AWS in platforms.", 0),
    ("The market and the developers disagree. Job postings cluster on Java and "
     "JavaScript, pay peaks on Swift, and usage peaks on JavaScript — popularity, "
     "hiring demand and salary are three different rankings.", 0),
    ("Demographics bound every claim: a young, degree-holding, "
     "US-and-Europe-weighted, self-selected sample.", 0),
], size=12)
fit(tb.text_frame)

# ------------------------------------------- 19. OVERALL FINDINGS & IMPLICATIONS
s = S[17]
ph = {p.placeholder_format.idx: p for p in s.placeholders}
write(ph[1].text_frame, [
    ("Findings", 0),
    ("", 0),
    (f"Job demand is concentrated: {facts['job_top3_share']}% of "
     f"{facts['job_total']:,} postings sit in three metro areas.", 0),
    (f"Hiring demand does not track developer usage: Java ({facts['skills_top'][1][1]:,} "
     f"postings) and JavaScript ({facts['skills_top'][2][1]:,}) outrank Python "
     f"({facts['skills_top'][3][1]:,}). The C total ({facts['skills_top'][0][1]:,}) is "
     "treated as unreliable — see A3.", 0),
    (f"Pay is led by {facts['sal_top'][0]} (${facts['sal_top'][1]:,}) — a "
     f"${facts['sal_spread']:,} spread over ten languages.", 0),
    ("PostgreSQL, TypeScript, Rust, Go and Redis gain share; MySQL, PHP, Java, "
     "SQL Server and jQuery lose it.", 0),
    (f"Location beats experience on pay: median compensation runs "
     f"${facts['comp_country_top3'][0][1]:,} in the "
     f"{facts['comp_country_top3'][0][0]} against "
     f"${facts['comp_country_bottom'][1][1]:,} in "
     f"{facts['comp_country_bottom'][1][0]}, while experience correlates only "
     f"r = {facts['corr_pairs']['comp_workexp']}.", 0),
    (f"Satisfaction is high ({facts['jobsat_7plus']}% score 7+) and effectively "
     f"uncorrelated with pay (r = {facts['corr_pairs']['jobsat_comp']}).", 0),
    (f"{facts['remote_pct']['Remote']}% work fully remotely and "
     f"{facts['remote_pct']['Hybrid (some remote, some in-person)']}% hybrid — only "
     f"{facts['remote_pct']['In-person']}% are fully in-person.", 0),
], size=11)
write(ph[2].text_frame, [
    ("Implications", 0),
    ("", 0),
    ("Recruit where the postings are, but source talent remotely — 85% of "
     "respondents already work remote or hybrid, so the concentration of postings "
     "need not constrain the hiring pool.", 0),
    ("Separate the three signals when planning: what developers use, what "
     "employers post for, and what the market pays are different rankings and "
     "should drive different decisions.", 0),
    ("Standardise on the consolidation winners (PostgreSQL, TypeScript, "
     "React/Next.js, AWS) to reduce long-run hiring risk.", 0),
    ("Fund Rust and Go upskilling now, while interest exceeds supply and before "
     "salaries reprice.", 0),
    ("Do not use pay as a retention lever on its own — the pay/satisfaction "
     "correlation is near zero, so flexibility and role design matter more.", 0),
    ("Set compensation bands by market geography rather than by seniority alone.", 0),
    ("Re-run this analysis each survey cycle: a single year shows a gap, a "
     "sequence of years shows a trend.", 0),
], size=11)
for t in (ph[1].text_frame, ph[2].text_frame):
    fit(t)

# ---------------------------------------------------------------- 20. CONCLUSION
s = S[18]
ph = {p.placeholder_format.idx: p for p in s.placeholders}
drop(ph[1])
picture(s, "a5_comp_by_age.png", 0.45, 2.05, 4.35, 3.6)
tb = by_text(s, "Point 1")
tb.left, tb.top, tb.width, tb.height = (Inches(5.15), Inches(1.55),
                                        Inches(7.5), Inches(5.0))
write(tb.text_frame, [
    ("The developer technology landscape is consolidating, not fragmenting. "
     "The leaders of today are the leaders of next year; the movement is in how "
     "much share each holds.", 0),
    ("Three independent rankings drive different decisions. Usage (survey), "
     "hiring demand (job postings) and pay (salary data) disagree — planning on "
     "any one of them alone is a mistake.", 0),
    ("The current-versus-desired gap is the most actionable output. Rust, Go, "
     "TypeScript, PostgreSQL and Redis gain share; PHP, Java, MySQL and legacy "
     "commercial engines lose it. That gap is the training and hiring agenda.", 0),
    ("Geography, not seniority, is the largest measurable driver of pay in this "
     "dataset, and pay is a weak lever on satisfaction.", 0),
    ("Limitations — a self-selected sample, roughly half of compensation values "
     "missing, one survey year with no time series, and stated intent rather than "
     "observed adoption. Conclusions are directional.", 0),
    ("Next steps — repeat across survey years to convert gaps into trends, join "
     "to a larger postings dataset for demand validation, and segment by role and "
     "region before acting.", 0),
], size=12)
fit(tb.text_frame)

# ------------------------------------------------------------------ 21. APPENDIX
s = S[19]
ph = {p.placeholder_format.idx: p for p in s.placeholders}
drop(ph[1])
tb = by_text(s, "Include any relevant additional charts")
tb.left, tb.top, tb.width, tb.height = (Inches(0.75), Inches(1.60),
                                        Inches(11.9), Inches(5.0))
write(tb.text_frame, [
    ("Additional charts and tables produced during the analysis phase", 0),
    ("A1 — Net mind-share shift by programming language", 0),
    ("A2 — Net mind-share shift by database", 0),
    ("A3 — Job postings by required technology", 0),
    ("A4 — Compensation distribution, raw and log-transformed", 0),
    ("A5 — Compensation by country", 0),
    ("A6 — Correlation matrix of pay, experience and satisfaction", 0),
    ("A7 — Job satisfaction distribution", 0),
    ("A8 — Language adoption by country", 0),
    ("A9 — Work arrangement and desired databases", 0),
    ("A10 — Dataset and methodology reference", 0),
], size=14)
fit(tb.text_frame)

# ------------------------------------------------------- appendix chart slides
APPENDIX = [
    ("A1 — NET MIND-SHARE SHIFT BY PROGRAMMING LANGUAGE", "r4_lang_netchange.png",
     "Each side is expressed as a share of all selections made on its own question, "
     f"because respondents name {facts['Language_per_resp'][0]} languages they have used "
     f"but only {facts['Language_per_resp'][1]} they want. Subtracting raw counts would push "
     "every language negative; comparing shares isolates the genuine change in priority."),
    ("A2 — NET MIND-SHARE SHIFT BY DATABASE", "r4_db_netchange.png",
     "Same share-based normalisation as A1 "
     f"({facts['Database_per_resp'][0]} databases used vs "
     f"{facts['Database_per_resp'][1]} wanted per respondent). PostgreSQL and Redis are the "
     "only large gainers; MySQL, SQL Server, MariaDB and Oracle all decline."),
    ("A3 — JOB POSTINGS BY REQUIRED TECHNOLOGY", "a1_job_postings_by_skill.png",
     f"C dominates at {facts['skills_top'][0][1]:,} postings, far ahead of Java "
     f"({facts['skills_top'][1][1]:,}) and JavaScript ({facts['skills_top'][2][1]:,}). "
     "The 'C' count is almost certainly inflated by substring matching on job text, "
     "which is why hiring-demand conclusions in this report lean on the per-location "
     "data rather than this breakdown."),
    ("A4 — COMPENSATION DISTRIBUTION", "a2_comp_distribution.png",
     f"Median ${facts['comp_median']:,} against a mean of ${facts['comp_mean']:,} on "
     f"{facts['comp_valid']:,} valid responses — a strong right skew. The log1p transform "
     "on the right makes the distribution usable for correlation work; the left panel is "
     "clipped at the 99th percentile for readability."),
    ("A5 — COMPENSATION BY COUNTRY", "a3_comp_by_country.png",
     f"Median pay ranges from ${facts['comp_country_top3'][0][1]:,} in the "
     f"{facts['comp_country_top3'][0][0]} to ${facts['comp_country_bottom'][1][1]:,} in "
     f"{facts['comp_country_bottom'][1][0]} — a "
     f"{round(facts['comp_country_top3'][0][1] / facts['comp_country_bottom'][1][1], 1)}× "
     "spread that dwarfs the effect of experience. Outliers are hidden and the axis is "
     "capped at the 98.5th percentile."),
    ("A6 — CORRELATION MATRIX", "a4_correlation.png",
     f"On {facts['corr_pairs']['n']:,} complete cases: compensation correlates weakly with "
     f"work experience (r = {facts['corr_pairs']['comp_workexp']}) and professional coding "
     f"years (r = {facts['corr_pairs']['comp_yearspro']}); job satisfaction is essentially "
     f"independent of both pay (r = {facts['corr_pairs']['jobsat_comp']}) and experience "
     f"(r = {facts['corr_pairs']['jobsat_yearspro']}). Listwise deletion means this subset "
     "over-represents respondents who answered the pay question."),
    ("A7 — JOB SATISFACTION DISTRIBUTION", "a9_jobsat.png",
     f"{facts['jobsat_n']:,} respondents answered ({facts['jobsat_missing_pct']}% did not). "
     f"The distribution is left-skewed toward high scores: median {facts['jobsat_median']:.0f}, "
     f"mean {facts['jobsat_mean']}, and {facts['jobsat_7plus']}% scoring 7 or above. "
     "The non-response rate is high enough that this should be read as the view of "
     "respondents who chose to answer."),
    ("A8 — LANGUAGE ADOPTION BY COUNTRY", "a8_language_by_country.png",
     "Each cell is the percentage of that country's respondents who report using the "
     "language, so rows are comparable despite very different sample sizes. The top ten "
     "languages are near-universal across all ten countries — regional variation is much "
     "smaller than the overall usage differences between languages."),
]
for title, img, caption in APPENDIX:
    sa = add_blank(title, title_size=22)
    picture(sa, img, 1.05, 1.15, 11.2, 4.55)
    note(sa, caption, 1.05, 5.85, 11.2, size=11.5, height=1.05)
    new_positions.append((sa, None))

# A9 — two small charts side by side
sa = add_blank("A9 — WORK ARRANGEMENT AND DESIRED DATABASES", title_size=22)
picture(sa, "a6_remote_work.png", 0.55, 1.30, 6.0, 4.2)
picture(sa, "a7_db_desired_pie.png", 6.75, 1.30, 6.0, 4.2)
note(sa,
     f"Left: only {facts['remote_pct']['In-person']}% of respondents work fully "
     f"in-person; {facts['remote_pct']['Remote']}% are fully remote and "
     f"{facts['remote_pct']['Hybrid (some remote, some in-person)']}% hybrid, which is why "
     "the geographic concentration of job postings need not constrain the hiring pool. "
     "Right: composition of the top five databases respondents want next year — "
     "PostgreSQL alone is roughly as large as the next two combined.",
     0.75, 5.70, 11.8, size=11.5, height=1.2)
new_positions.append((sa, None))

# A10 — reference table slide
sa = add_blank("A10 — DATASET AND METHODOLOGY REFERENCE", title_size=22)
rows = [
    ("Primary dataset", f"Stack Overflow Developer Survey extract — "
                        f"{facts['rows']:,} rows × {facts['cols']} columns, "
                        f"{facts['countries']} countries"),
    ("Secondary datasets", f"Jobs API postings ({facts['job_total']:,} across "
                           f"{facts['job_n_loc']} metros; plus a per-technology "
                           f"breakdown) and web-scraped salaries for 10 languages"),
    ("Multi-select fields", "LanguageHaveWorkedWith / WantToWorkWith, Database…, "
                            "Platform…, Webframe… — split on ';' and exploded to one "
                            "row per selection"),
    ("Selections per respondent",
     f"Languages {facts['Language_per_resp'][0]} used vs "
     f"{facts['Language_per_resp'][1]} wanted; databases "
     f"{facts['Database_per_resp'][0]} vs {facts['Database_per_resp'][1]}"),
    ("Compensation coverage", f"{facts['comp_valid']:,} valid values "
                              f"({facts['comp_missing_pct']}% missing); median "
                              f"${facts['comp_median']:,}, mean ${facts['comp_mean']:,}"),
    ("Outlier rule", f"1.5 × IQR; upper bound ${facts['comp_iqr_upper']:,}, "
                     f"{facts['comp_outliers']:,} values "
                     f"({facts['comp_outlier_pct']}%) excluded from pay comparisons"),
    ("Duplicates", f"{facts['dupes_full_row']} duplicate rows, "
                   f"{facts['dupes_response_id']} duplicate ResponseIds"),
    ("Tools", "Python (pandas, NumPy), SQLite/SQL, Matplotlib, Seaborn, "
              "BeautifulSoup, Google Looker Studio"),
    ("Chart colour", "Palette validated for colour-vision deficiency: purple = "
                     "current year, orange = next year, teal = demographics, "
                     "red = negative change"),
    ("Key caveat", "Self-selected sample; stated intent, not observed adoption; "
                   "a single survey year with no time series"),
]
tbl_shape = sa.shapes.add_table(len(rows) + 1, 2, Inches(0.72), Inches(1.20),
                                Inches(11.9), Inches(5.5))
table = tbl_shape.table
table.columns[0].width = Inches(3.05)
table.columns[1].width = Inches(8.85)
hdr = ["Item", "Detail"]
for c, text in enumerate(hdr):
    cell = table.cell(0, c)
    cell.text = text
    p = cell.text_frame.paragraphs[0]
    p.runs[0].font.size = Pt(12)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
for r, (k, v) in enumerate(rows, start=1):
    for c, text in enumerate((k, v)):
        cell = table.cell(r, c)
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(10.5)
        p.runs[0].font.bold = c == 0
        p.runs[0].font.color.rgb = INK if c == 0 else INK2
new_positions.append((sa, None))

# ------------------------------------------------------------- re-order slides
# Appended slides land at the end in creation order. Everything except the
# data-quality slide is appendix material and already sits where it belongs;
# only that first appended slide has to move up behind METHODOLOGY (index 4).
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
data_quality_el = ids[len(ids) - len(new_positions)]
sldIdLst.remove(data_quality_el)
sldIdLst.insert(5, data_quality_el)

prs.save(OUT_PPTX)

print(f"saved {OUT_PPTX}")
print(f"slides: {len(Presentation(OUT_PPTX).slides)}")
