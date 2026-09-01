"""Geometric / content QA for the built deck.

Stands in for a visual check: flags shapes off the canvas, overlapping picture and
text boxes, text frames without shrink-on-overflow, leftover template
placeholder text, and text volumes that are likely to overflow their box.
"""

import os

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

HERE = os.path.dirname(os.path.abspath(__file__))
DECK = os.path.join(HERE, "Data Analyst Capstone Project Report.pptx")

prs = Presentation(DECK)
SW, SH = prs.slide_width / 914400, prs.slide_height / 914400

LEFTOVERS = [
    "point1", "point2", "point3", "point 1", "point 2", "sub point",
    "finding 1", "implication 1", "goes here", "summarize key", "learner name",
    "<", "in module 1 you have", "please present", "include any relevant",
    "screenshot of dashboard", "dashboard insight\nsummarize",
]

issues = []


def box(sh):
    if sh.left is None or sh.top is None:
        return None
    return (sh.left / 914400, sh.top / 914400,
            (sh.left + (sh.width or 0)) / 914400,
            (sh.top + (sh.height or 0)) / 914400)


def overlap(a, b):
    if a is None or b is None:
        return 0.0
    ox = max(0, min(a[2], b[2]) - max(a[0], b[0]))
    oy = max(0, min(a[3], b[3]) - max(a[1], b[1]))
    return ox * oy


def has_autofit(sh):
    bodyPr = sh.text_frame._txBody.find(qn("a:bodyPr"))
    return bodyPr is not None and bodyPr.find(qn("a:normAutofit")) is not None


def est_lines(text, width_in, pt):
    """Rough wrap estimate: ~1.85 chars per pt of width for a UI sans."""
    if pt <= 0:
        return 0
    chars_per_line = max(8, int(width_in * 96 / (pt * 0.54)))
    lines = 0
    for para in text.split("\n"):
        lines += max(1, -(-len(para) // chars_per_line))
    return lines


for i, slide in enumerate(prs.slides, 1):
    pics, texts = [], []
    for sh in slide.shapes:
        b = box(sh)
        if b is None:
            continue

        # off-canvas
        if b[0] < -0.05 or b[1] < -0.05 or b[2] > SW + 0.05 or b[3] > SH + 0.05:
            issues.append(f"slide {i}: {sh.name!r} outside canvas "
                          f"({b[0]:.2f},{b[1]:.2f})-({b[2]:.2f},{b[3]:.2f})")

        if "PICTURE" in str(sh.shape_type):
            pics.append((sh, b))
            continue

        if sh.has_text_frame and sh.text_frame.text.strip():
            t = sh.text_frame.text
            low = t.lower()
            for token in LEFTOVERS:
                if token in low:
                    issues.append(f"slide {i}: leftover template text {token!r} "
                                  f"in {sh.name!r}: {t[:70]!r}")
                    break
            texts.append((sh, b, t))

            # overflow estimate
            sizes = [r.font.size.pt for p in sh.text_frame.paragraphs
                     for r in p.runs if r.font.size]
            pt = max(sizes) if sizes else 18
            lines = est_lines(t, b[2] - b[0], pt)
            need = lines * pt * 1.30 / 72
            avail = b[3] - b[1]
            if need > avail * 1.12 and not has_autofit(sh):
                issues.append(f"slide {i}: {sh.name!r} may overflow "
                              f"(~{need:.2f}in of text in {avail:.2f}in, no autofit)")

    # picture vs text collisions
    for psh, pb in pics:
        for tsh, tb, _ in texts:
            a = overlap(pb, tb)
            if a > 0.25:
                issues.append(f"slide {i}: picture {psh.name!r} overlaps text "
                              f"{tsh.name!r} by {a:.2f} sq in")
    for a_i in range(len(pics)):
        for b_i in range(a_i + 1, len(pics)):
            a = overlap(pics[a_i][1], pics[b_i][1])
            if a > 0.05:
                issues.append(f"slide {i}: pictures {pics[a_i][0].name!r} / "
                              f"{pics[b_i][0].name!r} overlap by {a:.2f} sq in")

print(f"slides: {len(prs.slides)}   canvas: {SW:.2f} x {SH:.2f} in")
if issues:
    print(f"\n{len(issues)} issue(s):")
    for x in issues:
        print("  -", x)
else:
    print("\nno geometric or leftover-text issues found")
